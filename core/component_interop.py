"""Singh360 component catalog publishing + PowerPoint interoperability helpers.

This module is intentionally additive. It works with ``LibraryV2`` without
requiring the component assets themselves to be committed to Git until the user
runs Publish Active Library.
"""
from __future__ import annotations

import html
import json
import re
import shutil
import tempfile
import uuid
from collections import defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

REPO_OWNER = "dstogsdill1"
REPO_NAME = "Singh360_SmartDraw"
REPO_BRANCH = "main"
PUBLIC_ROOT = "docs/component-library"


def _now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _stamp() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")


def _slug(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "-", (value or "component").strip()).strip("-") or "component"


def _read_json(path: Path, default: Any) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default


def _write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def component_is_archived(component: dict[str, Any]) -> bool:
    status = str(component.get("status") or "").strip().lower().replace(" ", "_")
    return bool(component.get("retired")) or status in {"retired", "archive", "archived", "deleted"}


def archive_component(lib: Any, comp_id: str) -> dict[str, Any]:
    return lib.update_component(
        comp_id,
        {
            "retired": True,
            "status": "archive",
            "approved": False,
            "needsReview": False,
        },
    )


def restore_component(lib: Any, comp_id: str) -> dict[str, Any]:
    return lib.update_component(
        comp_id,
        {
            "retired": False,
            "status": "active",
            "needsReview": False,
        },
    )


def permanent_delete_component(lib: Any, comp_id: str) -> dict[str, Any]:
    """Remove a component from active metadata and move unique assets to archive.

    Files are never destroyed without first being moved under ``.docs/archive``.
    Shared asset files remain in place when another component still references
    them.
    """
    payload = lib.load(include_legacy=True, include_retired=True)
    components = payload.get("components") or []
    target = next((c for c in components if c.get("id") == comp_id), None)
    if target is None:
        return {"ok": False, "error": "Component not found."}

    archive_dir = lib.archive / f"component_permanent_delete_{_stamp()}" / _slug(comp_id)
    archive_dir.mkdir(parents=True, exist_ok=True)
    _write_json(archive_dir / "component.json", target)

    # Back up both metadata sources before mutation.
    for p in (lib.manifest_path, lib._builder_export_path()):
        if p.exists():
            shutil.copy2(p, archive_dir / p.name)

    reference_keys = ("sourceFile", "edgeFile", "bwFile", "symbolFile", "thumbnailFile")
    used_by_others: set[str] = set()
    for c in components:
        if c.get("id") == comp_id:
            continue
        for key in reference_keys:
            rel = str(c.get(key) or "").replace("\\", "/")
            if rel:
                used_by_others.add(rel)

    moved: list[str] = []
    kept_shared: list[str] = []
    for key in reference_keys:
        rel = str(target.get(key) or "").replace("\\", "/")
        if not rel:
            continue
        if rel in used_by_others:
            kept_shared.append(rel)
            continue
        src = (lib.root / rel).resolve()
        if not src.is_file() or lib.root.resolve() not in src.parents:
            continue
        dst = archive_dir / "assets" / rel
        dst.parent.mkdir(parents=True, exist_ok=True)
        try:
            shutil.move(str(src), str(dst))
            moved.append(rel)
        except OSError:
            pass

    manifest = _read_json(lib.manifest_path, {"version": 2, "components": []})
    manifest["components"] = [c for c in manifest.get("components", []) if c.get("id") != comp_id]
    manifest["updatedAt"] = _now()
    _write_json(lib.manifest_path, manifest)

    export_path = lib._builder_export_path()
    if export_path.exists():
        export = _read_json(export_path, {"components": []})
        if isinstance(export, list):
            export = [c for c in export if c.get("id") != comp_id]
        elif isinstance(export, dict):
            export["components"] = [c for c in export.get("components", []) if c.get("id") != comp_id]
            export["updatedAt"] = _now()
        _write_json(export_path, export)

    return {
        "ok": True,
        "id": comp_id,
        "archiveDir": str(archive_dir),
        "assetsMoved": moved,
        "sharedAssetsKept": kept_shared,
    }


def _asset_path(lib: Any, component: dict[str, Any], key: str) -> Path | None:
    rel = str(component.get(key) or "").replace("\\", "/")
    if not rel:
        return None
    path = (lib.root / rel).resolve()
    if path.is_file() and lib.root.resolve() in path.parents:
        return path
    return None


def _copy_variant(lib: Any, component: dict[str, Any], key: str, out_dir: Path, base: str) -> str:
    src = _asset_path(lib, component, key)
    if src is None:
        return ""
    out_dir.mkdir(parents=True, exist_ok=True)
    dest = out_dir / f"{base}{src.suffix.lower()}"
    shutil.copy2(src, dest)
    return dest.as_posix()


def _github_raw(relative_to_repo: str) -> str:
    return f"https://raw.githubusercontent.com/{REPO_OWNER}/{REPO_NAME}/{REPO_BRANCH}/{relative_to_repo}"


def _public_catalog_html() -> str:
    return r'''<!doctype html>
<html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Singh360 Published Component Library</title>
<style>
body{margin:0;font:14px Arial;background:#f3f5f8;color:#111827}header{position:sticky;top:0;background:#0c1730;color:#fff;padding:14px 18px;z-index:5}h1{margin:0;font-size:20px}.tools{display:flex;gap:8px;padding:10px 18px;background:#fff;border-bottom:1px solid #ccd3df;position:sticky;top:51px;z-index:4}.tools input{flex:1;padding:8px}.tools select,.tools button{padding:8px}.grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(190px,1fr));gap:12px;padding:16px}.card{background:#fff;border:1px solid #d5dbe6;border-radius:8px;padding:9px}.img{height:135px;display:grid;place-items:center;background:#fafbfe;border:1px solid #e4e8ef}.img img{max-width:100%;max-height:100%;object-fit:contain}.name{font-weight:700;margin-top:7px}.meta{font-size:12px;color:#667085;margin:4px 0}.actions{display:flex;gap:5px;flex-wrap:wrap}.actions button,.actions a{font:12px Arial;padding:5px 7px;border:1px solid #aeb8c8;background:#fff;border-radius:4px;color:#111;text-decoration:none}.hidden{display:none}
</style></head><body><header><h1>Singh360 Published Component Library</h1></header>
<div class="tools"><input id="q" placeholder="Search components"><select id="cat"><option value="all">All categories</option></select><select id="rep"><option value="real">Real</option><option value="edge">Edge</option></select></div><div id="grid" class="grid"></div>
<script>
let items=[];const $=s=>document.querySelector(s);const esc=s=>String(s??'').replace(/[&<>"']/g,c=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]));
async function copyImage(url){const r=await fetch(url);const b=await r.blob();const img=await createImageBitmap(b);const c=document.createElement('canvas');c.width=img.width;c.height=img.height;c.getContext('2d').drawImage(img,0,0);const png=await new Promise(ok=>c.toBlob(ok,'image/png'));await navigator.clipboard.write([new ClipboardItem({'image/png':png})]);}
function render(){const q=$('#q').value.toLowerCase(),cat=$('#cat').value,rep=$('#rep').value;const f=items.filter(x=>(cat==='all'||x.category===cat)&&(!q||[x.displayName,x.partNumber,x.defaultLabel,x.category,(x.aliases||[]).join(' ')].join(' ').toLowerCase().includes(q)));$('#grid').innerHTML=f.map(x=>{const url=rep==='edge'?(x.edge||x.real):x.real;return `<article class="card"><div class="img">${url?`<img src="${esc(url)}">`:'No image'}</div><div class="name">${esc(x.displayName)}</div><div class="meta">${esc(x.category)}${x.partNumber?' · '+esc(x.partNumber):''}</div><div class="actions">${url?`<button onclick="copyImage('${esc(url)}')">Copy Image</button><a download href="${esc(url)}">Download</a>`:''}${x.real?`<a href="${esc(x.realRaw)}" target="_blank">GitHub Real</a>`:''}${x.edge?`<a href="${esc(x.edgeRaw)}" target="_blank">GitHub Edge</a>`:''}</div></article>`}).join('')}
fetch('catalog.json').then(r=>r.json()).then(d=>{items=d.components||[];const cats=[...new Set(items.map(x=>x.category))].sort();$('#cat').innerHTML+=[...cats].map(x=>`<option>${esc(x)}</option>`).join('');render()});$('#q').oninput=render;$('#cat').onchange=render;$('#rep').onchange=render;
</script></body></html>'''


def _raster_for_ppt(path: Path, temp_dir: Path) -> Path | None:
    ext = path.suffix.lower()
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}:
        return path
    out = temp_dir / f"{_slug(path.stem)}-{uuid.uuid4().hex[:6]}.png"
    try:
        if ext == ".svg":
            import cairosvg
            cairosvg.svg2png(url=str(path), write_to=str(out), output_width=1000, output_height=700)
            return out
        if ext == ".pdf":
            import fitz
            doc = fitz.open(str(path))
            pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=True)
            pix.save(str(out))
            doc.close()
            return out
    except Exception:
        return None
    return None


def build_powerpoint_template(output: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    output.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(17)
    prs.slide_height = Inches(11)
    blank = prs.slide_layouts[6]

    def add_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(blank)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(17), Inches(0.45))
        band.fill.solid(); band.fill.fore_color.rgb = RGBColor(38, 42, 48); band.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.25), Inches(0.07), Inches(16.5), Inches(0.28))
        p = tb.text_frame.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.65), Inches(16.5), Inches(8.35))
        body.fill.background(); body.line.color.rgb = RGBColor(135, 145, 160); body.line.width = Pt(1)
        note = slide.shapes.add_textbox(Inches(0.45), Inches(0.75), Inches(6.0), Inches(0.35))
        note.text_frame.text = subtitle
        note.text_frame.paragraphs[0].font.size = Pt(9); note.text_frame.paragraphs[0].font.color.rgb = RGBColor(100,105,115)
        safe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(9.18), Inches(16.5), Inches(1.5))
        safe.fill.solid(); safe.fill.fore_color.rgb = RGBColor(242, 244, 247); safe.line.color.rgb = RGBColor(170,175,185)
        txt = slide.shapes.add_textbox(Inches(0.6), Inches(9.62), Inches(15.8), Inches(0.5))
        txt.text_frame.text = "TITLE BLOCK / APP OUTPUT AREA — KEEP LAYOUT OBJECTS ABOVE THIS LINE"
        txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; txt.text_frame.paragraphs[0].font.size = Pt(11); txt.text_frame.paragraphs[0].font.bold = True; txt.text_frame.paragraphs[0].font.color.rgb = RGBColor(100,105,115)

    add_slide("SINGH360 EMS LAYOUT — BLANK", "17 × 11 landscape; use the body area for rough component placement.")
    add_slide("INTERIOR DEVICE LOCATION", "Use Singh360 component palette objects; preserve standard sizes.")
    add_slide("WICP LAYOUT", "Place safety symbols, WICP controls, notes, and signage callouts.")
    add_slide("RDM / IDF NETWORK LAYOUT", "Place RDM Data Manager, IDF/MDF, network devices, and paths.")
    while len(prs.slides) > 4:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId); del prs.slides._sldIdLst[-1]
    prs.save(output)
    return output


def build_powerpoint_palette(lib: Any, output: Path, variant: str = "real") -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    data = lib.load(include_legacy=True, include_retired=True)
    components = [c for c in data.get("components", []) if not component_is_archived(c)]
    by_cat: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for c in components:
        by_cat[str(c.get("category") or "custom")].append(c)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(); prs.slide_width = Inches(17); prs.slide_height = Inches(11); blank = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory(prefix="s360_ppt_") as tmp:
        tmpdir = Path(tmp)
        for category in sorted(by_cat):
            items = sorted(by_cat[category], key=lambda x: str(x.get("displayName") or "").lower())
            per_slide = 15
            for chunk_start in range(0, len(items), per_slide):
                chunk = items[chunk_start:chunk_start + per_slide]
                slide = prs.slides.add_slide(blank)
                band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(17), Inches(0.5))
                band.fill.solid(); band.fill.fore_color.rgb = RGBColor(12,23,48); band.line.fill.background()
                title = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(16.4), Inches(0.3))
                p = title.text_frame.paragraphs[0]; p.text = f"{category.replace('_',' ').title()} — {variant.title()} Components"; p.alignment = PP_ALIGN.CENTER; p.font.name="Arial"; p.font.size=Pt(16); p.font.bold=True; p.font.color.rgb=RGBColor(255,255,255)
                cols, rows = 5, 3; cell_w, cell_h = 3.28, 3.35
                for idx, comp in enumerate(chunk):
                    col, row = idx % cols, idx // cols
                    x, y = 0.25 + col * cell_w, 0.65 + row * cell_h
                    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cell_w - 0.12), Inches(cell_h - 0.14))
                    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(250,251,253); card.line.color.rgb=RGBColor(215,220,229); card.line.width=Pt(0.7)
                    key = "edgeFile" if variant == "edge" else "sourceFile"
                    src = _asset_path(lib, comp, key) or _asset_path(lib, comp, "sourceFile")
                    raster = _raster_for_ppt(src, tmpdir) if src else None
                    if raster:
                        dw = float(comp.get("defaultWidth") or 120) / 96.0
                        dh = float(comp.get("defaultHeight") or 64) / 96.0
                        max_w, max_h = cell_w - 0.55, cell_h - 1.05
                        scale = min(max_w / max(dw, 0.01), max_h / max(dh, 0.01), 1.0 if max(dw,dh)>1 else 3.0)
                        pw, ph = max(0.22, dw * scale), max(0.22, dh * scale)
                        pic = slide.shapes.add_picture(str(raster), Inches(x + (cell_w - 0.12 - pw)/2), Inches(y + 0.25 + (max_h - ph)/2), width=Inches(pw), height=Inches(ph))
                        pic.name = f"S360_COMPONENT_ID={comp.get('id')}|VARIANT={variant}"
                    label = slide.shapes.add_textbox(Inches(x+0.14), Inches(y+cell_h-0.82), Inches(cell_w-0.4), Inches(0.55))
                    tf=label.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=str(comp.get("displayName") or "Component"); p.alignment=PP_ALIGN.CENTER; p.font.name="Arial"; p.font.size=Pt(9); p.font.bold=True
                    part=str(comp.get("partNumber") or "")
                    if part:
                        p2=tf.add_paragraph(); p2.text=part; p2.alignment=PP_ALIGN.CENTER; p2.font.name="Arial"; p2.font.size=Pt(7); p2.font.color.rgb=RGBColor(95,100,110)
    while len(prs.slides) and len(prs.slides) > sum((len(v)+14)//15 for v in by_cat.values()):
        rId=prs.slides._sldIdLst[0].rId; prs.part.drop_rel(rId); del prs.slides._sldIdLst[0]
    prs.save(output)
    return output


def publish_active_library(lib: Any, repo_root: Path) -> dict[str, Any]:
    repo_root = Path(repo_root)
    target = repo_root / PUBLIC_ROOT
    if target.exists():
        backup = lib.archive / f"published_component_library_{_stamp()}"
        backup.parent.mkdir(parents=True, exist_ok=True)
        shutil.copytree(target, backup, dirs_exist_ok=True)
        shutil.rmtree(target)
    target.mkdir(parents=True, exist_ok=True)
    assets_root = target / "assets"; assets_root.mkdir(parents=True, exist_ok=True)

    data = lib.load(include_legacy=True, include_retired=True)
    published: list[dict[str, Any]] = []
    mapping: dict[str, Any] = {}
    for comp in data.get("components", []):
        if component_is_archived(comp):
            continue
        cid = str(comp.get("id") or _slug(str(comp.get("displayName") or "component")))
        cat = _slug(str(comp.get("category") or "custom"))
        folder = assets_root / cat / _slug(cid)
        real_abs = _copy_variant(lib, comp, "sourceFile", folder, "real")
        edge_abs = _copy_variant(lib, comp, "edgeFile", folder, "edge")
        if not real_abs and not edge_abs:
            continue
        real_rel = Path(real_abs).relative_to(repo_root).as_posix() if real_abs else ""
        edge_rel = Path(edge_abs).relative_to(repo_root).as_posix() if edge_abs else ""
        item = {
            "id": cid,
            "displayName": comp.get("displayName") or cid,
            "category": comp.get("category") or "custom",
            "collection": comp.get("collection") or "",
            "manufacturer": comp.get("manufacturer") or "",
            "partNumber": comp.get("partNumber") or "",
            "defaultLabel": comp.get("defaultLabel") or "",
            "aliases": comp.get("aliases") or [],
            "tags": comp.get("tags") or [],
            "defaultWidth": comp.get("defaultWidth") or 120,
            "defaultHeight": comp.get("defaultHeight") or 64,
            "real": Path(real_rel).relative_to(PUBLIC_ROOT).as_posix() if real_rel else "",
            "edge": Path(edge_rel).relative_to(PUBLIC_ROOT).as_posix() if edge_rel else "",
            "realRaw": _github_raw(real_rel) if real_rel else "",
            "edgeRaw": _github_raw(edge_rel) if edge_rel else "",
        }
        published.append(item)
        mapping[cid] = {"real": item["realRaw"], "edge": item["edgeRaw"]}

    catalog = {"version": 1, "generatedAt": _now(), "components": published}
    _write_json(target / "catalog.json", catalog)
    (target / "index.html").write_text(_public_catalog_html(), encoding="utf-8")
    (target / "README.md").write_text(
        "# Singh360 Published Component Library\n\nGenerated from active local components.\n\n"
        f"GitHub Pages URL (after Pages is enabled from `/docs`):\n\n"
        f"https://{REPO_OWNER}.github.io/{REPO_NAME}/component-library/\n",
        encoding="utf-8",
    )
    _write_json(lib.root / "published_map.json", {"version": 1, "generatedAt": _now(), "components": mapping})
    build_powerpoint_template(target / "Singh360_EMS_17x11_Layout_Template.pptx")
    build_powerpoint_palette(lib, target / "Singh360_Component_Library_Real.pptx", "real")
    build_powerpoint_palette(lib, target / "Singh360_Component_Library_Edge.pptx", "edge")
    return {
        "ok": True,
        "published": len(published),
        "folder": str(target),
        "localUrl": "/published-components/",
        "githubPagesUrl": f"https://{REPO_OWNER}.github.io/{REPO_NAME}/component-library/",
    }


def read_published_map(lib: Any) -> dict[str, Any]:
    return _read_json(lib.root / "published_map.json", {"version": 1, "components": {}})
