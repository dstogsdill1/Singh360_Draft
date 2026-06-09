"""extractors/panel_config.py — control-panel configuration (.pptx / .pdf).

Panel Configuration decks/PDFs show what's inside each control panel (WICP):
breakers, contactors, controllers. We pull the readable text (pptx slides via
python-pptx; pdf via PyMuPDF) and create one PANEL node per slide/page, with
the recognized panel name and the raw text kept as an attribute for review.
Heavy libs are imported lazily; if unavailable we flag rather than fail.
"""
from __future__ import annotations

import re
from pathlib import Path

from core.model import ProjectModel, Node, NodeKind, slug

_PANEL_RE = re.compile(r"\b(WICP\s*\d+|PANEL\s+[A-Z0-9\-]+|[A-Z]{2,4}CP\d*)\b", re.I)


def _slide_texts_pptx(path: Path):
    try:
        from pptx import Presentation
    except ImportError:
        return None
    out = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
        out.append((i + 1, "\n".join(chunks)))
    return out


def _page_texts_pdf(path: Path):
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return None
    out = []
    doc = fitz.open(path)
    for i in range(doc.page_count):
        out.append((i + 1, doc[i].get_text("text").strip()))
    doc.close()
    return out


def extract(path: str | Path, model: ProjectModel) -> None:
    path = Path(path)
    model.note_source(str(path))

    if path.suffix.lower() in (".pptx", ".ppt"):
        pages = _slide_texts_pptx(path)
        if pages is None:
            model.flag("review", f"{path.name}: install python-pptx to read panel decks", path.name)
            return
    elif path.suffix.lower() == ".pdf":
        pages = _page_texts_pdf(path)
        if pages is None:
            model.flag("review", f"{path.name}: install PyMuPDF to read panel PDFs", path.name)
            return
    else:
        model.flag("review", f"{path.name}: unsupported panel-config type", path.name)
        return

    n = 0
    for page_no, text in pages or []:
        if not text:
            continue
        m = _PANEL_RE.search(text)
        pname = (m.group(1).strip() if m else f"Panel (slide {page_no})")
        model.add_node(Node(
            id=slug("panel", pname, str(page_no)),
            kind=NodeKind.PANEL, name=pname,
            attrs={"text": text[:1500]}, source=f"{path.name}:p{page_no}",
        ))
        n += 1
    model.flag("info", f"Panel config: {n} panels indexed from {path.name}", path.name)
